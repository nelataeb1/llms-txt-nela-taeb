"""Builds docs/llms-txt-deck.pptx, the interview deck for this project.

Requirements: pip install python-pptx
Usage:        python scripts/build_deck.py

Theme is lifted from the reference deck (Google Slides export): 070707
background, 0F0F0F cards, 333333 hairlines, Helvetica Neue / Helvetica Neue
Light, F3F3F3 titles over 888888 body copy, dashed frame rules and a top-left
eyebrow + title header.

On-slide text stays sparse; the narration lives in the speaker notes.
Architecture diagrams are the dark-theme Mermaid renders in
docs/deck/diagrams (regenerate with scripts/render_diagrams.sh).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DIAGRAMS = DOCS / "deck" / "diagrams"
OUT = DOCS / "llms-txt-deck.pptx"
LOGO = DOCS / "deck" / "profound-logo.png"

CANDIDATE = "Nela Taeb"
LIVE_URL = "https://llms-txt-nela-taeb.vercel.app"

BG = RGBColor(0x07, 0x07, 0x07)
CARD = RGBColor(0x0F, 0x0F, 0x0F)
HAIRLINE = RGBColor(0x33, 0x33, 0x33)
TITLE = RGBColor(0xF3, 0xF3, 0xF3)
BRIGHT = RGBColor(0xF5, 0xF5, 0xF0)
MUTED = RGBColor(0x88, 0x88, 0x88)
LABEL = RGBColor(0xC8, 0xC8, 0xC0)

SANS = "Helvetica Neue"
SANS_LIGHT = "Helvetica Neue Light"
MONO = "Roboto Mono"


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #
def text_box(slide, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    frame.paragraphs[0].alignment = align
    return frame


def write(
    frame,
    text,
    *,
    size,
    color=MUTED,
    font=SANS_LIGHT,
    bold=False,
    italic=False,
    space_before=0,
    line_spacing=1.25,
    first=False,
    align=None,
):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.line_spacing = line_spacing
    para.space_before = Pt(space_before)
    if align is not None:
        para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return para


def line(slide, x1, y1, x2, y2, *, dashed=False, arrow=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = HAIRLINE
    conn.line.width = Emu(9525)
    if dashed:
        conn.line.dash_style = 4  # long dash, as in the reference deck
    if arrow:
        from lxml import etree

        ln = conn.line._get_or_add_ln()
        tail = etree.SubElement(
            ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd"
        )
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
    return conn


def card(slide, x, y, w, h, *, fill=CARD, border=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = HAIRLINE
        shape.line.width = Emu(10125)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def caption(slide, x, y, w, text, *, align=PP_ALIGN.CENTER, size=7.0):
    frame = text_box(slide, x, y, w, 0.18, align=align)
    write(frame, text, size=size, color=LABEL, first=True, align=align)


def picture(slide, path: Path, x, y, w, *, max_h=None, border=True):
    """Places a screenshot at width ``w``, shrinking it to respect ``max_h``."""
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if max_h is not None and pic.height > Inches(max_h):
        scale = Inches(max_h) / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = Inches(x) + int((Inches(w) - pic.width) / 2)
    if border:
        pic.line.color.rgb = HAIRLINE
        pic.line.width = Emu(10125)
    return pic


def bottom(shape) -> float:
    return (shape.top + shape.height) / 914400


def diagram(slide, name: str, x, y, max_w, max_h):
    """Places a Mermaid render, scaled to fit the box and horizontally centred."""
    path = DIAGRAMS / f"{name}.png"
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y))
    scale = min(Inches(max_w) / pic.width, Inches(max_h) / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = Inches(x) + int((Inches(max_w) - pic.width) / 2)
    pic.top = Inches(y) + int((Inches(max_h) - pic.height) / 2)
    return pic


def notes(slide, text: str):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text.strip()


# --------------------------------------------------------------------------- #
# slide chrome
# --------------------------------------------------------------------------- #
def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def content_slide(prs, eyebrow: str, title: str):
    slide = new_slide(prs)
    line(slide, 0.23, 0.0, 0.23, 5.62, dashed=True)
    line(slide, 9.75, 0.0, 9.75, 5.62, dashed=True)
    line(slide, 0.0, 0.32, 10.0, 0.32, dashed=True)

    frame = text_box(slide, 0.7, 0.58, 8.6, 0.22)
    write(frame, eyebrow, size=8.5, color=MUTED, font=SANS, first=True)
    frame = text_box(slide, 0.7, 0.83, 8.8, 0.74)
    write(frame, title, size=30, color=TITLE, font=SANS_LIGHT, first=True, line_spacing=1.0)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.19), Inches(5.07), width=Inches(1.34))
    return slide


def wrapped_height(text: str, width: float, size: float) -> float:
    chars_per_line = max(int(width * 72 / (0.52 * size)), 12)
    lines = 1
    used = 0
    for word in text.split():
        need = len(word) + (1 if used else 0)
        if used + need > chars_per_line:
            lines += 1
            used = len(word)
        else:
            used += need
    return lines * size * 1.34 / 72


def rows(slide, items, *, x=0.75, y=1.62, w=8.5, gap=0.3, title_size=11.5, body_size=10.0):
    top = y
    for head, body in items:
        frame = text_box(slide, x, top, w, 0.22)
        write(frame, head, size=title_size, color=TITLE, first=True)
        top += title_size * 1.34 / 72 + 0.09
        if body:
            height = wrapped_height(body, w, body_size)
            frame = text_box(slide, x, top, w, height)
            write(frame, body, size=body_size, color=MUTED, first=True, line_spacing=1.3)
            top += height
        top += gap
    return top


def row_pair(slide, left_item, right_item, *, y, title_size=10.0, body_size=8.5):
    """Two single-line rows side by side — used under a full-width diagram."""
    for (head, body), x in ((left_item, 0.7), (right_item, 5.1)):
        frame = text_box(slide, x, y, 4.2, 0.22)
        write(frame, head, size=title_size, color=TITLE, first=True)
        frame = text_box(slide, x, y + title_size * 1.34 / 72 + 0.09, 4.2, 0.2)
        write(frame, body, size=body_size, color=MUTED, first=True)


def numbered(slide, items, *, y=1.75, gap=0.38, w=7.8, head_size=13.0, body_size=10.5):
    top = y
    for index, (head, body) in enumerate(items):
        frame = text_box(slide, 0.79, top + 0.02, 0.4, 0.28)
        write(frame, f"{index + 1:02d}", size=11, color=TITLE, font=SANS, first=True)
        frame = text_box(slide, 1.29, top, w, 0.32)
        write(frame, head, size=head_size, color=TITLE, first=True)
        top += head_size * 1.34 / 72 + 0.1
        if body:
            height = wrapped_height(body, w, body_size)
            frame = text_box(slide, 1.29, top, w, height)
            write(frame, body, size=body_size, color=MUTED, first=True, line_spacing=1.3)
            top += height
        top += gap
    return top


def card_grid(slide, items, *, x=0.75, y=1.7, w=4.2, h=1.34, gapx=0.24, gapy=0.19, cols=2):
    for index, (head, body) in enumerate(items):
        col, row = index % cols, index // cols
        left = x + col * (w + gapx)
        top = y + row * (h + gapy)
        card(slide, left, top, w, h)
        frame = text_box(slide, left + 0.22, top + 0.22, w - 0.44, 0.22)
        write(frame, head, size=11.5, color=TITLE, first=True)
        if body:
            frame = text_box(slide, left + 0.22, top + 0.48, w - 0.44, h - 0.68)
            write(frame, body, size=9.5, color=MUTED, first=True, line_spacing=1.3)


def code_block(slide, x, y, w, lines_, *, size=8.5):
    height = 0.28 + len(lines_) * (size + 5) / 72
    card(slide, x, y, w, height)
    frame = text_box(slide, x + 0.18, y + 0.14, w - 0.36, height - 0.28)
    for index, text in enumerate(lines_):
        write(frame, text, size=size, color=LABEL, font=MONO, first=index == 0, line_spacing=1.25)
    return height


def tradeoff(slide, text, *, y=4.66, x=0.7, w=8.6):
    frame = text_box(slide, x, y, w, 0.3)
    para = write(frame, "Trade-off  ", size=9, color=TITLE, font=SANS, first=True)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.name = SANS_LIGHT
    run.font.color.rgb = MUTED


def chips(slide, labels, *, x=0.7, y=1.5, gap=0.16, size=9.0, height=0.3):
    """A single row of small pill labels — the deck's sparse alternative to bullets."""
    left = x
    for text in labels:
        width = 0.26 + len(text) * size * 0.0092
        card(slide, left, y, width, height)
        frame = text_box(slide, left, y, width, height, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        write(frame, text, size=size, color=LABEL, first=True, align=PP_ALIGN.CENTER)
        left += width + gap
    return left


# --------------------------------------------------------------------------- #
# slides
# --------------------------------------------------------------------------- #
def slide_01_title(prs):
    slide = new_slide(prs)
    line(slide, 0.37, 0.06, 0.37, 5.31, dashed=True)
    line(slide, 9.59, 0.06, 9.59, 5.31, dashed=True)
    line(slide, 0.0, 5.33, 10.0, 5.33, dashed=True)

    frame = text_box(slide, 0.55, 1.09, 8.9, 0.7)
    write(frame, "Automated llms.txt Generator", size=42, color=TITLE, font=SANS, first=True, line_spacing=1.0)
    frame = text_box(slide, 0.62, 1.9, 8.9, 0.5)
    write(frame, "Making a site legible to answer engines", size=28, color=TITLE, first=True, line_spacing=1.0)
    frame = text_box(slide, 0.64, 2.68, 8.6, 0.3)
    write(frame, LIVE_URL, size=11, color=LABEL, font=MONO, first=True)

    frame = text_box(slide, 0.78, 4.2, 5.75, 0.6)
    write(frame, f"{CANDIDATE}  |  FDE Candidate", size=16, color=TITLE, italic=True, first=True)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(7.24), Inches(4.2), width=Inches(2.11))

    notes(
        slide,
        """
        Quick framing before I open anything: the take-home asked for a tool that generates an
        llms.txt file for a domain and keeps it current. What I built is a deployed Next.js app
        that takes a URL, crawls the site, and produces a spec-valid llms.txt, then monitors that
        site and re-publishes the file when the content actually changes.

        I also went past the spec in two places, because a file that is merely valid is not
        necessarily useful: an AI-readiness audit that grades how the domain looks to real answer
        engine crawlers, and a retrieval eval that measures whether the generated index can
        actually route a question to the right page.

        The app is live at the URL on this slide — everything I show is running in production on
        Vercel with Neon Postgres behind it.
        """,
    )


def slide_02_agenda(prs):
    slide = content_slide(prs, "Agenda", "What I'll walk through")
    numbered(
        slide,
        [
            ("Problem & the standard", ""),
            ("Product overview and the core pipeline", ""),
            ("Beyond spec: audit and retrieval eval", ""),
            ("Live demo", ""),
            ("Technical deep dive, component by component", ""),
            ("Trade-offs and what I'd build next", ""),
        ],
        y=1.62,
        gap=0.12,
        head_size=13.5,
    )
    notes(
        slide,
        """
        Structure: I'll spend the first few minutes on why llms.txt matters and what the spec
        actually requires, then show the product and the pipeline end to end. After the demo I go
        component by component — crawl, extract, rank, group, render, validate, monitor, persist —
        with the trade-off I made in each one. I'll close on what I would build next if this were
        a real product rather than a take-home.

        If you want to interrupt at any point, please do — the deep-dive slides are independent,
        so we can spend the time wherever it's most useful for you.
        """,
    )


def slide_03_problem(prs):
    slide = content_slide(prs, "The Problem", "AI answers are the new front page")
    card_grid(
        slide,
        [
            ("Discovery moved", "Answer engines reply instead of returning ten links."),
            ("HTML is hostile to agents", "Nav, banners and JS shells bury the content."),
            ("Context budgets are small", "A crawler reads a fraction of a site, if any."),
            ("Hand-authoring doesn't scale", "Someone must curate — and re-curate on every change."),
        ],
        y=1.75,
        w=4.2,
        h=1.4,
        gapy=0.22,
    )
    notes(
        slide,
        """
        The reason this file exists at all: for a growing share of queries, nobody sees a results
        page any more. ChatGPT, Perplexity and Google's AI surfaces read a handful of sources and
        synthesise an answer, and if the model can't cheaply read your site, your brand is simply
        absent from that answer. That's exactly the visibility problem Profound sells into.

        Why is a website hard to read? A marketing page is mostly navigation, cookie banners,
        tracking scripts and client-rendered shells. An agent with a limited context window burns
        most of it on markup before it reaches the paragraph that mattered — and on a
        client-rendered SPA it may extract nothing at all.

        llms.txt is the proposed fix, and the analogy is robots.txt: one file at a well-known path
        that tells an agent where to look. The difference is that robots.txt is exclusion and
        llms.txt is curation — you have to choose the right pages and describe them well.

        And that curation is the actual work. A docs site with a few hundred pages needs someone
        to pick which ones matter, write a one-line summary for each, and then keep that in sync
        every time docs, pricing or product pages ship. That's what I automated.
        """,
    )


def slide_04_standard(prs):
    slide = content_slide(prs, "The Standard", "What llms.txt has to look like")
    rows(
        slide,
        [
            ("H1 — required", "one title, first line"),
            ("Blockquote — summary", "what the site is, in a sentence"),
            ("H2 sections of link lists", "[title](url): note"),
            ("Optional — the escape hatch", "safe to skip under context pressure"),
        ],
        x=0.7,
        y=1.65,
        w=4.1,
        gap=0.24,
        title_size=11,
        body_size=9.5,
    )
    code_block(
        slide,
        5.1,
        1.65,
        4.2,
        [
            "# Profound",
            "",
            "> AI search visibility platform.",
            "",
            "## Documentation",
            "- [Quickstart](/docs/quickstart): first crawl",
            "- [API](/docs/api): REST endpoints",
            "",
            "## Optional",
            "- [Privacy](/legal/privacy): policy",
        ],
        size=8.0,
    )
    caption(slide, 5.1, 4.55, 4.2, "markdown, /llms.txt, one H1 — validated on every render", align=PP_ALIGN.LEFT)
    notes(
        slide,
        """
        The spec at llmstxt.org is deliberately small, which is why it's easy to produce something
        that looks right and is actually invalid. The shape: exactly one H1 with the site name; an
        optional blockquote that summarises the site; then free-form detail; then H2 sections,
        each of which is a markdown list of links, each link optionally followed by a colon and a
        short note.

        Two details drive most of my design. First, the notes matter more than the links — the
        note is the only signal an agent has for choosing between two plausible pages, so a bad
        description is worse than a missing page. Second, the trailing "Optional" section is the
        spec's explicit escape hatch: anything an agent can safely skip when it's short on context
        goes there. I use it for legal pages and for the overflow beyond my link budget.

        There's also llms-full.txt, which inlines page content rather than linking to it. I
        generate it on request, but it's a different trade: complete, but far too big for most
        context windows.
        """,
    )


def slide_05_overview(prs):
    slide = content_slide(prs, "Product Overview", "Paste a URL, get a maintained file")
    rows(
        slide,
        [
            ("Generate", "crawl → rank → sections → llms.txt"),
            ("Validate", "graded against the spec, in the UI"),
            ("Monitor", "daily re-crawl, page-level diffs"),
            ("Serve", "/s/{siteId}/llms.txt, text/plain"),
        ],
        x=0.7,
        y=1.72,
        w=4.0,
        gap=0.26,
        title_size=11,
        body_size=9.5,
    )
    pic = picture(slide, DOCS / "generate.png", 5.05, 1.68, 4.25, max_h=2.9)
    caption(
        slide,
        5.05,
        bottom(pic) + 0.12,
        4.25,
        "tryprofound.com — crawled, grouped, spec compliant",
        align=PP_ALIGN.LEFT,
    )
    notes(
        slide,
        """
        The product surface is deliberately one input. You paste a domain, and within about a
        minute you have a spec-valid llms.txt, a validation report next to it, and — if you tick
        monitoring — a permanent URL that always serves the current version of that file.

        Four verbs. Generate is the pipeline I'll show next. Validate re-parses the rendered text
        and grades it, so the score describes the exact bytes a consumer would download rather
        than the builder's intentions. Monitor re-crawls on a daily cron and diffs page by page.
        Serve exposes the latest snapshot at a stable URL, as text/plain, so a customer can proxy
        it from their own domain without redeploying anything.

        Worth saying explicitly: no database and no API key are required to run it. Without
        DATABASE_URL it falls back to an in-memory store, and without an OpenAI key the
        heuristics produce the whole file. The LLM only improves wording.
        """,
    )


def slide_06_pipeline(prs):
    slide = content_slide(prs, "How it works", "The core pipeline")
    diagram(slide, "pipeline", 0.5, 1.7, 9.0, 2.4)
    frame = text_box(slide, 0.7, 4.35, 8.6, 0.5)
    write(
        frame,
        "Every stage is a pure function over serializable state.",
        size=11,
        color=TITLE,
        first=True,
    )
    write(frame, "A crawl can pause, persist and resume across serverless requests.", size=9.5, color=MUTED, space_before=4)
    notes(
        slide,
        """
        Left to right: a URL comes in; the crawler resolves robots.txt and sitemaps and collects
        pages; extraction pulls title, description, canonical, links and a content hash out of
        each one; classification assigns a page kind and a score; grouping turns the ranked list
        into sections under a link budget; rendering emits the markdown; validation re-parses it.

        The important property is the one written at the bottom: each stage is a pure function
        from serializable state to serializable state. That isn't architectural purity for its
        own sake — it's what makes the crawl survive a serverless platform. Any stage can stop
        mid-way, be written to Postgres, and picked up by a different Lambda invocation.

        LLM enrichment hangs off the side deliberately. It sees only the pages the crawler already
        found and can rewrite titles, section names and descriptions — it never contributes a URL,
        so it cannot hallucinate a page into the file. If the key is missing, the deterministic
        path still produces a valid document.
        """,
    )


def slide_07_beyond_spec(prs):
    slide = content_slide(prs, "Beyond the spec", "Two features the brief didn't ask for")
    card_grid(
        slide,
        [
            (
                "AI-readiness audit",
                "Six real bot probes, cloaking detection, weighted checks → an A–F grade with a fix per failure.",
            ),
            (
                "Retrieval eval",
                "Two LLMs, eight questions: can an agent holding only the file reach the right page?",
            ),
        ],
        y=1.8,
        w=4.2,
        h=1.75,
    )
    chips(slide, ["valid ≠ reachable", "valid ≠ useful", "both are numbers, not opinions"], y=3.9)
    notes(
        slide,
        """
        The brief asked for generation plus automated updates. I added two things because the
        obvious question after "is the file valid?" is "does it do anything?", and validity
        answers neither of the two ways this fails in the real world.

        First failure mode: the file is perfect but the crawler never gets to read it — robots.txt
        blocks GPTBot, or a WAF serves a challenge page to anything that isn't a browser. The
        audit tests that directly by fetching the site as each named bot and comparing what came
        back with a browser-like baseline, so it catches cloaking that a normal crawl never sees.

        Second failure mode: the file is fetched and parsed, but the descriptions are so generic
        that the agent picks the wrong link. The retrieval eval quantifies that with a
        writer/reader pair — I'll show the mechanics later.

        Both were also chosen because they're the questions Profound's customers ask: why am I
        invisible in AI answers, and what do I change first.
        """,
    )


def slide_08_demo(prs):
    slide = content_slide(prs, "Live Demo", "What I'll show")
    rows(
        slide,
        [
            ("1 — Generate", "paste a domain, watch slices progress"),
            ("2 — Inspect", "sections, notes, validation report"),
            ("3 — Monitor", "track, refresh, page-level diff"),
            ("4 — Audit + eval", "grade the domain, score the file"),
        ],
        x=0.7,
        y=1.75,
        w=3.9,
        gap=0.3,
        title_size=11,
        body_size=9.5,
    )
    pic = picture(slide, DOCS / "monitored-sites.png", 4.95, 1.75, 4.35, max_h=2.6)
    caption(
        slide,
        4.95,
        bottom(pic) + 0.12,
        4.35,
        "tracked sites, last check, served file URL",
        align=PP_ALIGN.LEFT,
    )
    frame = text_box(slide, 0.7, 4.7, 8.6, 0.3)
    write(frame, LIVE_URL, size=11, color=BRIGHT, font=MONO, first=True)
    notes(
        slide,
        """
        Demo script, in case I lose the network and have to talk through screenshots.

        One: paste tryprofound.com and hit generate. Progress ticks in slices rather than sitting
        on a spinner, because each slice is its own request — that's the serverless constraint
        showing through in the UX, and it's a feature: you can see how far it got.

        Two: the result. Sections come from page kind, the ordering comes from the ranking score,
        and each link has a note. Next to it, the validation panel grades the exact rendered text.

        Three: tick monitoring and hit refresh. If nothing changed, you get "no changes" and no
        snapshot is written; if something changed, you get a page-level diff — added, removed,
        edited — and a new snapshot. The file URL underneath is stable across all of that.

        Four: run the audit and the retrieval eval on the same result, which is the natural
        segue into the two deep-dive slides later.
        """,
    )


def slide_09_architecture(prs):
    slide = content_slide(prs, "Technical Deep Dive", "System architecture")
    diagram(slide, "system", 0.5, 1.5, 9.0, 3.3)
    caption(slide, 0.7, 4.95, 8.6, "Next.js 15 · TypeScript · Vercel functions + cron · Neon Postgres · Vitest", align=PP_ALIGN.LEFT)
    notes(
        slide,
        """
        The whole thing is one Next.js 15 app on Vercel. Route handlers are thin: parse and
        validate input with Zod, call a library function, return JSON. All the logic lives in
        src/lib as plain TypeScript functions with no framework or platform imports, which is why
        it's straightforward to unit test — the test suite never boots a server.

        Three entry points into the same pipeline. /api/generate starts and advances a job.
        /api/sites tracks, refreshes and lists history. /api/audit and /api/eval run the two
        beyond-spec features against an existing result. The cron endpoint calls the same
        refreshSite() the manual button does, so there's exactly one refresh code path.

        Persistence is behind a Store interface with two implementations: Neon Postgres in
        production, in-memory for local and for CI. And GET /s/{siteId}/llms.txt reads the latest
        snapshot straight from the store, so serving is a database read rather than a crawl.

        The dotted edge is the OpenAI call — optional, and only ever enrichment.
        """,
    )


def slide_10_crawl(prs):
    slide = content_slide(prs, "Deep Dive — Crawl", "Slice-based, resumable crawling")
    diagram(slide, "crawl", 0.5, 1.45, 9.0, 2.5)
    row_pair(
        slide,
        ("Sitemap first, links as fallback", "links still supply nav + inbound signals"),
        ("Polite and bounded", "crawl-delay ≤ 2s, concurrency 6, 120 pages"),
        y=4.1,
    )
    tradeoff(
        slide,
        "more state to serialize, but the crawl never dies at the serverless timeout.",
    )
    notes(
        slide,
        """
        The constraint that shaped this component: a Vercel function has a hard execution limit,
        and crawling a hundred pages politely takes longer than that. The usual answer is a
        queue and a worker; I didn't want to add that infrastructure for a take-home, so I made
        the crawl itself interruptible.

        initCrawl resolves robots.txt, collects sitemap URLs and returns a state object — frontier,
        visited set, pages, inbound-link counts, nav URLs, counters. crawlSlice takes that state
        and a millisecond budget, fetches batches of six concurrently until the budget or the page
        cap is hit, and returns new state. The route persists it and returns progress; the browser
        calls again. So a crawl is a loop of ordinary HTTP requests, each one comfortably inside
        the timeout.

        Sitemaps first because they give breadth for one or two fetches, including sitemap indexes
        and gzipped sitemaps. Link crawling still runs, partly as a fallback for sites without a
        usable sitemap, and partly because I need the link graph: inbound counts and navigation
        membership are two of the strongest ranking signals I have.

        Politeness is not optional if you're going to point this at other people's sites: robots
        is respected by default, crawl-delay is honoured but capped at two seconds so one hostile
        value can't hang a run, and concurrency is six.
        """,
    )


def slide_11_extract(prs):
    slide = content_slide(prs, "Deep Dive — Extract", "Turning HTML into a page record")
    card_grid(
        slide,
        [
            ("Title", "h1 → og:title → <title>, site suffix stripped"),
            ("Description", "meta → og → twitter → JSON-LD → first sentences"),
            ("Identity", "canonical URL, normalized; sha content hash"),
            ("Signals", "nav membership, inbound links, word count, markdown alternate"),
        ],
        y=1.66,
        w=4.2,
        h=1.22,
        gapy=0.16,
    )
    code_block(
        slide,
        0.75,
        4.45,
        8.5,
        ['<link rel="alternate" type="text/markdown" href="/docs/api.md">   →   preferred link target'],
        size=8.0,
    )
    notes(
        slide,
        """
        Extraction is Cheerio, not a browser, and it's a cascade of fallbacks per field because
        every site is wrong in a different way. Title: h1, then og:title, then the title tag, with
        the trailing " | Site Name" suffix stripped so section lists don't read like a stutter.
        Description: meta description, then Open Graph, then Twitter cards, then JSON-LD, and only
        then the first couple of sentences of body text.

        Two things here are load-bearing later. The canonical URL is resolved and normalized —
        trailing slashes, index files, tracking params, fragment stripping — because otherwise the
        same page shows up three times in the file and the duplicate-page check in the audit
        misfires. And every page gets a sha content hash over the fields that matter, which is
        what makes monitoring cheap: a refresh compares hashes rather than diffing text.

        The markdown alternate is a small detail with a big payoff. If a site publishes a
        text/markdown alternate — increasingly common in docs tooling — I link to that instead of
        the HTML page, so the agent gets clean prose with no markup to strip.
        """,
    )


def slide_12_rank(prs):
    slide = content_slide(prs, "Deep Dive — Rank & Group", "Which pages earn a place")
    diagram(slide, "rank", 0.5, 1.4, 9.0, 2.55)
    row_pair(
        slide,
        ("Readable formula, not a model", "every placement is explainable to a customer"),
        ("Budget forces a choice", "overflow and legal fall into Optional"),
        y=4.1,
    )
    tradeoff(
        slide,
        "hand-tuned weights are transparent and free; embeddings would generalise "
        "better at a model call per page.",
    )
    notes(
        slide,
        """
        A crawl of a real site returns far more pages than belong in an index, so something has to
        choose. scorePage is deliberately a readable formula rather than a model: page kind sets
        the base — home 100, docs 40, api 38, guides 34, product 26, legal 2 — then depth costs 8
        per level, inbound links add 2.5 each up to a cap, being in the navigation adds 18,
        appearing in the sitemap, having a markdown alternate and having a real description all
        add, and content volume adds a little with penalties for pagination-style paths.

        I can defend every one of those numbers in an interview, which is exactly the point: when
        a customer asks why their pricing page ranked below the API reference, I can answer from
        the formula instead of shrugging at a model.

        Grouping then walks the ranked list into fixed-order sections by page kind — Overview,
        Documentation, API Reference, Guides, and so on — under a link budget of 25 per section
        and 120 overall. Oversized "Other" buckets get split by top-level path so you don't get one
        section with sixty links. Legal pages and everything past the budget fall into Optional,
        which is the spec's own escape hatch for agents under context pressure.

        The honest weakness: the weights encode my assumptions about how sites are laid out. They
        do well on docs-shaped and marketing-shaped sites and less well on, say, an e-commerce
        catalogue. Clustering embeddings would generalise, at the cost of a model call per page
        and a much harder answer to "why is this page here".
        """,
    )


def slide_13_render(prs):
    slide = content_slide(prs, "Deep Dive — Render", "From ranked sections to markdown")
    rows(
        slide,
        [
            ("Deterministic document model", "buildDocument() → renderLlmsTxt(): one shape, one writer"),
            ("Notes are the product", "extracted description, trimmed to a single clause"),
            ("llms-full.txt on request", "same structure, page text inlined"),
        ],
        x=0.7,
        y=1.65,
        w=4.15,
        gap=0.28,
        title_size=10.5,
        body_size=9.0,
    )
    code_block(
        slide,
        5.15,
        1.65,
        4.15,
        [
            "## API Reference",
            "- [REST API](/docs/api): endpoints,",
            "  auth and rate limits",
            "- [Webhooks](/docs/webhooks): event",
            "  payloads and retries",
        ],
        size=8.0,
    )
    tradeoff(
        slide,
        "heuristic notes are instant and safe; LLM notes read better but need a key — "
        "and never add a URL.",
    )
    notes(
        slide,
        """
        Rendering is split in two on purpose: buildDocument produces a typed document — title,
        summary, sections, links, notes — and renderLlmsTxt turns that into markdown. Nothing else
        in the codebase writes llms.txt syntax, so there's one place where the format lives, and
        the validator can be written against the output without sharing code with the builder.

        The notes are where the quality actually is. A link line without a note is nearly useless
        to an agent choosing between two pages, so each note comes from the best description I
        could extract, trimmed to a single clause. When an OpenAI key is present, gpt-4o-mini
        rewrites titles, section names and notes in one batched call — but it's handed the list of
        pages and asked only to relabel them. It never returns URLs, which is the guardrail that
        makes the optional model safe to enable.

        llms-full.txt uses the same document and inlines page text. It's genuinely useful for a
        small docs site and completely impractical for a large one, so it's opt-in.
        """,
    )


def slide_14_validate(prs):
    slide = content_slide(prs, "Deep Dive — Validate", "Grading the exact bytes we emit")
    card_grid(
        slide,
        [
            ("Re-parse, don't trust", "the checker reads the rendered text, not the builder's state"),
            ("Errors", "missing/duplicate H1, content before title, malformed links, H3+"),
            ("Warnings", "duplicate URLs, missing summary, links without notes"),
            ("One function, three callers", "generator UI, POST /api/validate, the audit"),
        ],
        y=1.72,
        w=4.2,
        h=1.35,
        gapy=0.2,
    )
    tradeoff(
        slide,
        "re-parsing costs milliseconds and buys an independent check on the bytes we ship.",
    )
    notes(
        slide,
        """
        validateLlmsTxt takes a string and returns errors and warnings. The deliberate choice is
        that it re-parses the rendered document rather than inspecting the in-memory model. If the
        renderer has a bug, a state-based check would happily agree with it; a text-based check
        catches it. That's also what lets me point the same function at somebody else's live
        /llms.txt during the audit, where I have no model at all — only bytes.

        The severity split follows the spec: structural violations are errors — no H1, more than
        one H1, content before the title, link lines that don't parse, headings deeper than H2.
        Things that are legal but degrade usefulness are warnings — the same URL listed twice, a
        missing summary blockquote, links with no note.

        Three callers share it: the generator UI grades every result inline, POST /api/validate
        lets you paste an arbitrary file, and the audit uses it on whatever the domain is already
        serving. One definition of "valid", no drift.
        """,
    )


def slide_15_monitor(prs):
    slide = content_slide(prs, "Deep Dive — Monitor", "Keeping the file current")
    diagram(slide, "monitor", 0.5, 1.5, 9.0, 2.9)
    caption(
        slide,
        0.7,
        4.3,
        8.6,
        "one refreshSite() behind both the daily cron and the manual button",
        align=PP_ALIGN.LEFT,
    )
    tradeoff(
        slide,
        "hash diffing makes 'nothing changed' nearly free, but a formatting-only edit is "
        "invisible to it.",
    )
    notes(
        slide,
        """
        Automated updates were an explicit requirement, and the naive version — re-render nightly
        and store the result — produces a history where every day looks like a change and nobody
        can see what actually happened.

        So the unit of comparison is the page, not the file. Every page carries a content hash;
        diffPages compares the previous snapshot's pages with the new crawl and reports added,
        removed and changed URLs, including what moved in the title or description. Those page
        hashes roll up into a snapshot hash. If the snapshot hash is unchanged, I update
        lastCheckedAt and write nothing. If it changed, I store a new snapshot with the diff
        attached.

        Two consequences. History stays legible — the timeline only contains days where something
        really happened — and storage stays small, because I drop page text from stored snapshots
        unless llms-full.txt needs it.

        Both triggers call the same function, which matters more than it sounds: the daily Vercel
        cron and the Refresh button in the UI cannot drift apart. The cost of the hash approach is
        that a change which doesn't affect the hashed fields — pure formatting, say — reads as no
        change at all. For this use case that's the right side of the trade.
        """,
    )


def slide_16_data(prs):
    slide = content_slide(prs, "Deep Dive — Persistence", "Two stores, one interface")
    diagram(slide, "data", 0.5, 1.75, 5.3, 1.4)
    rows(
        slide,
        [
            ("Zero-config default", "no DATABASE_URL → in-memory; the app still runs end to end"),
            ("Schema on first use", "no migration step to run before a deploy"),
            ("Jobs are state, not threads", "a crawl in progress is a row, so any instance can advance it"),
        ],
        x=6.0,
        y=1.75,
        w=3.3,
        gap=0.24,
        title_size=10,
        body_size=8.5,
    )
    tradeoff(
        slide,
        "one interface keeps local runs and CI dependency-free, at the cost of writing "
        "every query twice.",
    )
    notes(
        slide,
        """
        Three tables. Sites: the tracked domain, its crawl options, whether monitoring is on, and
        when it was last checked. Snapshots: the rendered llms.txt, its content hash, the page
        records and the diff, one row per actual change. Jobs: an in-flight crawl — the serialized
        crawl state, progress, and the finished result.

        That jobs table is what makes the slice-based crawl work on serverless. The crawl isn't a
        thread anybody has to keep alive; it's a row. Whichever function instance handles the next
        request loads the state, advances it by one slice, and writes it back.

        Everything goes through a Store interface with a Neon Postgres implementation and an
        in-memory one. It means you can clone the repo and run the app with no infrastructure at
        all, and the whole test suite runs without a database. The price is honest: two
        implementations to keep in sync, and I can't reach for Postgres-specific features like
        JSONB indexes or LISTEN/NOTIFY without breaking the in-memory path.

        Schema is created on first use rather than through a migration tool — right call for a
        take-home, first thing I'd replace with real migrations in production.
        """,
    )


def slide_17_audit(prs):
    slide = content_slide(prs, "Beyond Spec — Audit", "How an answer engine sees the domain")
    diagram(slide, "audit", 0.5, 1.4, 9.0, 2.6)
    row_pair(
        slide,
        ("Weighted score → A–F grade", "every failing check ships the concrete fix"),
        ("Use case", "why a brand is invisible, and what to change first"),
        y=4.1,
    )
    tradeoff(
        slide,
        "live bot probes cost seconds and politeness budget, so the audit is opt-in.",
    )
    notes(
        slide,
        """
        The audit answers "can an answer engine reach this content at all", and it does it by
        measurement rather than by linting.

        For each of six named agents — GPTBot, OAI-SearchBot, ChatGPT-User, ClaudeBot,
        PerplexityBot and Google-Extended — I evaluate the site's robots.txt with that exact
        user-agent, then actually fetch the entry page identifying as that bot. The response is
        compared with a browser-like baseline fetch. If GPTBot gets a fraction of the bytes a
        browser gets, that's a WAF challenge page or a JS stub, and it's the single most common
        reason a brand is missing from AI answers. You cannot detect that by reading robots.txt.

        Around that sit content checks: sitemap coverage, missing descriptions, thin and duplicate
        pages, markdown alternates, canonical tags, whether content is server-rendered — plus, if
        the domain already publishes an /llms.txt, I fetch it, validate it, and compare it against
        my fresh crawl to find stale links and pages the live file never mentions.

        The checks are weighted into a score and a grade, A down to F, and every failure carries
        the specific fix. That's the deliverable a customer can act on: not "you scored 62" but
        "GPTBot is blocked by this line in robots.txt, and 40 of your docs pages have no
        description".
        """,
    )


def slide_18_eval(prs):
    slide = content_slide(prs, "Beyond Spec — Eval", "Is the file actually usable?")
    diagram(slide, "eval", 0.5, 1.4, 9.0, 2.6)
    row_pair(
        slide,
        ("Ids cross the wire, not URLs", "a hallucinated answer misses instead of scoring"),
        ("Misses are actionable", "names the note that failed to disambiguate"),
        y=4.1,
    )
    tradeoff(
        slide,
        "two model calls per run cost money and vary slightly, but eight questions still "
        "catch vague notes.",
    )
    notes(
        slide,
        """
        This is the one I'd defend hardest as a product idea. A file can be perfectly valid and
        still useless, because the notes don't distinguish between pages. The eval measures that.

        A writer model sees the crawled pages with their content and writes eight questions a real
        user might ask, each tagged with the link id that should answer it. A reader model then
        sees only the generated llms.txt — no page content — and answers each question by choosing
        a link id. Accuracy is how often the reader lands on the page the writer had in mind.

        The design detail that makes it trustworthy: both sides exchange numbered ids, never URLs.
        If the reader invents something, the id doesn't resolve and it's a miss — it can't get
        silent credit for a plausible-looking URL.

        And the failures are the useful output. Every miss names the link that was chosen instead
        and why its note didn't disambiguate, which is precisely the description to rewrite. Run
        it before and after an edit and you can see whether the file got better, which is what
        turns this from a vibe into a regression test.
        """,
    )


def slide_19_tradeoffs(prs):
    slide = content_slide(prs, "Engineering", "Trade-offs I'd defend")
    numbered(
        slide,
        [
            ("Slices over one long request", "survives the timeout; costs serialized state"),
            ("Heuristics first, LLM optional", "always valid without a key; the model only relabels"),
            ("Hash diffing over re-render comparison", "cheap no-ops; misses formatting-only edits"),
            ("Audit and eval opt-in", "keeps generation fast; two extra clicks"),
            ("Interface over two stores", "zero-config and testable; every query written twice"),
        ],
        y=1.62,
        gap=0.16,
        head_size=12,
        body_size=9.0,
    )
    notes(
        slide,
        """
        Pulling the trade-offs into one place, because this is usually the part of the interview
        that matters.

        Slice crawling: I chose more moving parts over a crawl that dies at 60 seconds. The state
        has to be serializable and each slice is a round-trip, but progress is visible and nothing
        is lost when a function is recycled.

        Heuristics first: the deterministic path always produces a spec-valid file. The model is
        strictly an enhancement, and it's structurally prevented from inventing URLs. A missing
        key degrades wording, never correctness.

        Hash diffing: cheap no-ops and readable history, in exchange for being blind to changes
        that don't touch the hashed fields.

        Opt-in audit and eval: six bot probes plus two LLM calls would double the latency of every
        generation for a feature most runs don't need.

        Two store implementations: worth it for a project that has to be clonable and testable
        without infrastructure; the first thing I'd collapse if this became a real product.

        On testing, since it usually comes up: the suite is Vitest over the pure functions — robots
        and sitemap parsing including wildcards, anchors, indexes and gzip; extraction fallbacks
        and URL normalization; classification and ranking order; sectioning and link budgets;
        validation errors and warnings; snapshot diffing; and audit scoring. Plus lint, typecheck
        and build in the same command.
        """,
    )


def slide_20_next(prs):
    slide = content_slide(prs, "Roadmap", "What I'd build next")
    numbered(
        slide,
        [
            ("Headless rendering for JS-only sites", "fall back to a browser on thin pages only"),
            ("Embeddings for ranking and sectioning", "name sections from content, not URL paths"),
            ("Incremental recrawl", "ETag / If-Modified-Since, per-page scheduling"),
            ("Event-driven refresh", "sitemap lastmod or a deploy webhook — minutes, not hours"),
            ("Score trends over time", "audit and eval are both better as a time series"),
        ],
        y=1.55,
        gap=0.1,
        head_size=12,
        body_size=9.0,
    )
    frame = text_box(slide, 0.79, 4.6, 8.5, 0.4)
    write(frame, LIVE_URL, size=12, color=BRIGHT, font=MONO, first=True)
    write(frame, f"{CANDIDATE} — thank you", size=9.5, color=MUTED, space_before=3)
    notes(
        slide,
        """
        Five things, roughly in the order I'd do them.

        Headless rendering: the biggest correctness gap. A fully client-rendered site extracts as
        an empty page today. I'd keep Cheerio as the fast path and fall back to a real browser
        only for pages that come back suspiciously thin, so I pay for it on the few pages that
        need it.

        Embeddings: cluster page embeddings to derive sections from what pages are about rather
        than from URL segments, and use similarity to catch near-duplicates the hash misses. The
        cost is explainability, so I'd want the heuristic score kept alongside it.

        Incremental recrawl: right now a refresh re-crawls the whole site. ETag and
        If-Modified-Since plus per-page scheduling — hot docs pages hourly, legal pages monthly —
        would cut the daily cost by an order of magnitude on a large site.

        Event-driven refresh: polling a sitemap's lastmod, or better, a CMS or deploy webhook, so
        the file is minutes behind a publish instead of up to a day.

        And trends: both the audit grade and the retrieval accuracy are far more interesting as a
        line over time than as a single number — that's the dashboard I'd build if this were a
        product rather than a take-home.

        That's the tour — happy to go deeper anywhere, and the live app is on the screen if you'd
        like to throw a domain at it.
        """,
    )


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_fn in (
        slide_01_title,
        slide_02_agenda,
        slide_03_problem,
        slide_04_standard,
        slide_05_overview,
        slide_06_pipeline,
        slide_07_beyond_spec,
        slide_08_demo,
        slide_09_architecture,
        slide_10_crawl,
        slide_11_extract,
        slide_12_rank,
        slide_13_render,
        slide_14_validate,
        slide_15_monitor,
        slide_16_data,
        slide_17_audit,
        slide_18_eval,
        slide_19_tradeoffs,
        slide_20_next,
    ):
        slide_fn(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(f"wrote {build().relative_to(ROOT)}")
